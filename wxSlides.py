from collections import defaultdict
from pathlib import Path
import subprocess
import textwrap

import numpy as np

import wx
from wx.media import MediaCtrl, MC_NO_AUTORESIZE
from wx.grid import Grid

import pptx
from pptx.slide import Slide
from pptx.util import Pt, Cm, Inches, Length
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from lxml import etree

from PIL import Image, ImageDraw, ImageFont  # , ImageText
from moviepy import ImageClip, ColorClip, CompositeVideoClip


def xpath(el, query: str):
    """Utility to query an `pptx.shapes.Shape`'s xml tree."""
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    return etree.ElementBase.xpath(el, query, namespaces=nsmap)


def set_table_font_size(table, size_pt):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(size_pt)


def set_text_font_size(tbox, size_pt):
    for paragraph in tbox.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size_pt)


def move_slide(pres, from_index: int, to_index: int) -> None:
    """Move slide at position `from_index` in presentation `pres` to `to_index`"""
    slides = list(pres.slides._sldIdLst)
    if to_index < 0:
        to_index = len(slides) + to_index
    pres.slides._sldIdLst.remove(slides[from_index])

    pres.slides._sldIdLst.insert(to_index, slides[from_index])


def autoplay_media(media) -> None:
    """
    Utility to autoplay a media (currently only video) upon on entering the slide containing it.

    Args:
        media: `Shape` object containing the video.
    """
    el_id = xpath(media.element, ".//p:cNvPr")[0].attrib["id"]
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), ".//p:cond")[0]
    cond.set("delay", "0")
    cond.set("evt", "onBegin")


def get_thumbnail_from_video(movie_file: str, img_format: str = ".jpg") -> str:
    video_input_path = Path(movie_file).resolve()
    img_output_path = video_input_path.parent / (video_input_path.stem + img_format)
    subprocess.call(
        [
            "ffmpeg",
            "-i",
            video_input_path,
            "-ss",
            "00:00:00.000",
            "-vframes",
            "1",
            img_output_path,
            "-y",
        ],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return str(img_output_path)


def text_to_scrolling_video(
    text,
    video_file,
    /,
    *,
    width=454,
    height=56,
    font_size=20,
    hmargin=0,
    speed=50,
    fps=24,
):

    _width = width * 2
    _height = height * 2
    _font_size = font_size * 2
    _speed = speed * 2
    _hmargin = hmargin * 2

    # 1. Load the Pillow Font
    font = ImageFont.load_default(size=_font_size)

    # 2. Wrap the text dynamically
    # We calculate the exact pixel width of a sample alphabet to find average char width
    avg_char_width = font.getlength(text) / len(text)
    _hmargin += avg_char_width
    max_chars = int((_width - (2 * _hmargin)) / avg_char_width)
    wrapped_text = "\n".join(textwrap.wrap(text, width=max_chars - 2))

    # 3. Calculate the required height of the final image
    # We use a dummy image to measure exactly how tall the text block will be
    dummy_draw = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    bbox = dummy_draw.multiline_textbbox((0, 0), wrapped_text, font=font)
    text_height = bbox[3] - bbox[1]
    text_width = bbox[2] - bbox[0]

    # Add a little buffer to the bottom just to be safe
    img_height = text_height + 2 * _font_size
    # img_width = text_width + 2 * _hmargin

    # 4. Generate the Text Image Canvas (Transparent Background)
    # We make the image exactly the width of the video, so X positioning is locked to 0
    img = Image.new("RGBA", (_width, img_height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    # Draw the text! We apply the hmargin physically inside the image
    draw.multiline_text(
        (_hmargin, 0),
        wrapped_text,
        font=font,
        fill="black",
        align="left",
        # spacing=avg_char_width,
    )

    # 5. Send directly to MoviePy
    # Convert Pillow Image to a NumPy array, then into a MoviePy ImageClip
    img_array = np.asarray(img)
    txt_clip = ImageClip(img_array)

    # 6. Calculate duration and Animate
    delta_h = img_height - _height

    # We only move the Y coordinate. X stays 0 because margins are baked into the image.
    txt_clip = (
        txt_clip.with_position(lambda t: (0, -_speed * t)).with_duration(
            delta_h / _speed
        )
        if delta_h > 0
        else txt_clip.with_duration(1.0 / fps)
    )

    bg_clip = ColorClip(
        (_width, _height), color=(255, 255, 255), duration=txt_clip.duration
    )
    # 7. Composite and Save
    video = CompositeVideoClip([bg_clip, txt_clip]).without_audio().resized(0.5)
    video.write_videofile(
        f"{video_file}.mp4",
        fps=fps,
        # codec="libx264",
        # audio_codec="aac",
        # threads=8,
        # logger=None,
        # preset="veryfast",
    )
    video.save_frame(f"{video_file}.png")
    # print(wrapped_text)


def add_movie(
    pres,
    slide,
    movie_file: str,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    mime_type: str = "video/mp4",
    poster_frame_image: str | None = None,
    add_fullscreen: bool = True,
    hide_fullscreen_slide: bool = True,
):
    """
    Wrapper around add_movie method of a `pptx.slide.Slide` instance to add movies with functionality to toggle fullscreen mode

    Args:
        pres: the presentaion instance which contains the slide instance
        slide: slide instance to which we add the movie
        movie_file: path to the movie .mp4 file
        left: X-coordinate of the movie frame's top-left corner
        top: Y-coordinate of the movie frame's top-left corner
        width: width of the movie frame
        height: height of the movie frame
        mime_type: input to the mime_type keyword argument of slide.add_movie method.
        poster_frame_image: input to the poster_frame_image keyword argument of slide.add_movie method.
        add_fullscreen: Whether to add fullscreen toggling feature.
        hide_fullscreen_slide: Whether to hide the extra fullscreen slide or not. Recommend setting True if using PowerPoint and False if using Keynote.

    Returns:
        If `add_fullscreen == True`, returns a tuple of `(movie_shape, fullscreen_movie_slide)` else just returns the `movie_shape`
    """

    if add_fullscreen:
        thn_img_path = (
            poster_frame_image
            if poster_frame_image is not None
            else get_thumbnail_from_video(movie_file)
        )
        thn_img = slide.shapes.add_picture(thn_img_path, left, top, width, height)

        fs_btn_w, fs_btn_h = Inches(1.5), Inches(0.5)
        fs_btn_left = left + width - fs_btn_w - Inches(0.2)
        fs_btn_top = top + height - fs_btn_h - Inches(0.2)

        fs_btn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, fs_btn_left, fs_btn_top, fs_btn_w, fs_btn_h
        )
        fs_btn.text = "Fullscreen"

        fs_movie_slide = pres.slides.add_slide(pres.slide_layouts[6])
        if hide_fullscreen_slide:
            fs_movie_slide.element.set("show", "0")

        movie = fs_movie_slide.shapes.add_movie(
            movie_file,
            0,
            0,
            pres.slide_width,
            pres.slide_height,
            mime_type=mime_type,
            poster_frame_image=thn_img_path,
        )
        # movie.click_action.hyperlink.address = None
        autoplay_media(movie)
        fs_movie_slide.element.append(
            parse_xml(
                '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" advOnClick="0"/>'
            )
        )

        fs_btn.click_action.target_slide = fs_movie_slide

        fs_exit_btn = fs_movie_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            pres.slide_width - fs_btn_w - Inches(0.2),
            Inches(0.2),
            fs_btn_h,
            fs_btn_h,
        )
        fs_exit_btn.text = "X"
        fs_exit_btn.click_action.target_slide = slide

        return movie, fs_movie_slide, thn_img

    movie = slide.shapes.add_movie(
        movie_file,
        left,
        top,
        width,
        height,
        mime_type=mime_type,
        poster_frame_image=None,
    )
    return movie


class wxShape(wx.StaticBoxSizer):
    def __init__(self, parent, orient=wx.HORIZONTAL, title=""):
        super().__init__(orient, parent, label=title)


class wxTextBox(wxShape):
    def __init__(self, parent, orient=wx.HORIZONTAL, title="", **kwargs):
        super().__init__(parent, orient=orient, title=title)
        self.textCtrl = wx.TextCtrl(self.StaticBox, **kwargs)
        self.Add(self.textCtrl, wx.SizerFlags(1).Expand().Border())

    @property
    def Text(self):
        return self.textCtrl.GetValue()

    @Text.setter
    def Text(self, text):
        self.textCtrl.ChangeValue(text)

    def SaveToSlide(
        self,
        slide,
        left,
        top,
        width,
        height,
        scroll=False,
        wrap=True,
        font_size=None,
        **kwargs,
    ):
        if scroll:
            video_file_prefix = (
                f"./Videos/scrolling_text_{slide.slide_id}_{len(slide.shapes) + 1}"
            )
            text_to_scrolling_video(
                self.Text,
                video_file_prefix,
                width=int(width.pt),
                height=int(height.pt),
                font_size=font_size,
                **kwargs,
            )

            return slide.shapes.add_movie(
                f"{video_file_prefix}.mp4",
                left,
                top,
                width,
                height,
                mime_type="video/mp4",
                poster_frame_image=f"{video_file_prefix}.png",
            )

        tbox = slide.shapes.add_textbox(left, top, width, height)
        tbox.text_frame.text = self.Text
        tbox.text_frame.word_wrap = wrap
        if font_size is not None:
            set_text_font_size(tbox, font_size)
        return tbox


class wxMovieShape(wxShape):
    def __init__(
        self,
        parent,
        orient=wx.HORIZONTAL,
        title="",
        file_name=None,
        style=MC_NO_AUTORESIZE,
        **kwargs,
    ):
        super().__init__(parent, orient=orient, title=title)
        self.movieCtrl = MediaCtrl(self.StaticBox, style=style, **kwargs)
        self.movieCtrl.ShowPlayerControls()
        self.Add(self.movieCtrl, wx.SizerFlags(1).Expand().Border())
        if file_name is not None:
            self.LoadVideo(file_name)

    def LoadVideo(self, file_name, thumbnail_file_name=None):
        self.fileName = file_name
        self.thumbFileName = thumbnail_file_name
        status = self.movieCtrl.Load(self.fileName)
        if not status:
            raise ValueError(f"Can't load {self.fileName}, not a valid video file")

    def SaveToSlide(
        self,
        pres,
        slide,
        left,
        top,
        width,
        height,
        mime_type="video/mp4",
    ):
        return add_movie(
            pres,
            slide,
            self.fileName,
            left,
            top,
            width,
            height,
            mime_type,
            poster_frame_image=self.thumbFileName,
        )


class wxTableShape(wxShape):
    def __init__(self, parent, orient=wx.HORIZONTAL, title="", **kwargs):
        super().__init__(parent, orient=orient, title=title)
        self.tableCtrl = Grid(self.StaticBox, **kwargs)
        self.tableCtrl.UseNativeColHeader()
        self.Add(self.tableCtrl, wx.SizerFlags(1).Expand().Border())

    def CreateGrid(self, num_rows, num_cols):
        self.tableCtrl.CreateGrid(num_rows, num_cols)

    def SetColumnNames(self, column_names):
        numCols = self.tableCtrl.GetNumberCols()
        if len(column_names) != numCols:
            raise ValueError(
                f"Sequence of column_names (given length {len(column_names)}) must be of same length as num_cols ({numCols})"
            )
        for icol, colName in enumerate(column_names):
            self.tableCtrl.SetColLabelValue(icol, str(colName))

    def SetTableData(self, tabular_data):
        for irow, row_data in enumerate(tabular_data):
            for icol, cell_data in enumerate(row_data):
                self.tableCtrl.SetCellValue(irow, icol, str(cell_data))

    def SetFromDataframe(self, df):
        numRows, numCols = df.shape
        self.CreateGrid(numRows, numCols)
        self.SetColumnNames(df.columns)
        self.SetTableData(df.values.tolist())

    def SaveToSlide(self, slide, left, top, width, height, font_size=None):
        rows = self.tableCtrl.GetNumberRows()
        cols = self.tableCtrl.GetNumberCols()
        tableFrame = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
        table = tableFrame.table
        for icol in range(cols):
            table.cell(0, icol).text = self.tableCtrl.GetColLabelValue(icol)

        for irow in range(rows):
            for icol in range(cols):
                table.cell(irow + 1, icol).text = self.tableCtrl.GetCellValue(
                    irow, icol
                )

        if font_size is not None:
            set_table_font_size(table, font_size)

        return tableFrame


class wxSlide(wx.Panel):
    LAYOUT = 6  # BLANK SLIDE

    def __init__(self, parent, **kwargs):
        super().__init__(parent, **kwargs)
        self.shapes = defaultdict(list)
        self.SlideLayout()

    def MakeSlideLayout(self):
        return wx.BoxSizer(wx.VERTICAL)

    def SlideLayout(self):
        lSizer = self.MakeSlideLayout()
        mainSizer = self.GetSizer()

        if mainSizer is None:
            self.SetSizer(lSizer)
        else:
            mainSizer.Add(lSizer)
            self.Layout()

    def AddTextBox(self, *args, **kwargs):
        self.shapes["textbox"].append(wxTextBox(self, *args, **kwargs))

    def AddMovie(self, *args, **kwargs):
        self.shapes["movie"].append(wxMovieShape(self, *args, **kwargs))

    def AddTable(self, *args, **kwargs):
        self.shapes["table"].append(wxTableShape(self, *args, **kwargs))

    def SaveToPres(self, pres):
        # print(f"Slide: {len(pres.slides)}, layout: {self.LAYOUT}")
        return pres.slides.add_slide(pres.slide_layouts[self.LAYOUT])


class wxTitleSlide(wxSlide):
    LAYOUT = 0

    def MakeSlideLayout(self):
        mainSizer = super().MakeSlideLayout()
        self.AddTextBox(title="Title")
        self.AddTextBox(title="Subtitle")
        self.Title = "Example title"
        self.SubTitle = "Example subtitle"
        for shape in self.shapes["textbox"]:
            mainSizer.Add(shape, wx.SizerFlags(0).Expand().DoubleBorder())
        return mainSizer

    @property
    def Title(self):
        return self.shapes["textbox"][0].Text

    @Title.setter
    def Title(self, value):
        self.shapes["textbox"][0].Text = str(value)

    @property
    def SubTitle(self):
        return self.shapes["textbox"][1].Text

    @SubTitle.setter
    def SubTitle(self, value):
        self.shapes["textbox"][1].Text = str(value)

    def SaveToPres(self, pres):
        slide = super().SaveToPres(pres)
        slide.shapes.title.text = self.Title
        slide.placeholders[1].text = self.SubTitle
        return slide


class wxTitleOnlySlide(wxSlide):
    LAYOUT = 5

    def MakeSlideLayout(self):
        mainSizer = super().MakeSlideLayout()
        self.AddTextBox(title="Title")
        self.shapes["textbox"][0].Text = "Example title"
        mainSizer.Add(
            self.shapes["textbox"][0], wx.SizerFlags(0).Expand().DoubleBorder()
        )
        return mainSizer

    @property
    def Title(self):
        return self.shapes["textbox"][0].Text

    @Title.setter
    def Title(self, value):
        self.shapes["textbox"][0].Text = str(value)

    def SaveToPres(self, pres):
        slide = super().SaveToPres(pres)
        slide.shapes.title.text = self.Title
        return slide


class wxStepSlide(wxTitleOnlySlide):
    def __init__(
        self,
        parent,
        title,
        text,
        movie_file_name=None,
        movie_thumbnail_file_name=None,
        **kwargs,
    ):
        super().__init__(parent, **kwargs)
        self.Title = title
        self.shapes["textbox"][1].Text = text

        if movie_file_name is not None:
            self.shapes["movie"][0].LoadVideo(
                movie_file_name, movie_thumbnail_file_name
            )

    def MakeSlideLayout(self):
        mainSizer = super().MakeSlideLayout()

        self.AddTextBox(title="Text:", style=wx.TE_MULTILINE | wx.TE_BESTWRAP)
        self.AddMovie(title="Video:")

        # for key, val in self.shapes.items():
        #    print(key, len(val))

        topSizer = wx.BoxSizer(wx.HORIZONTAL)
        topSizer.Add(
            self.shapes["textbox"][1], wx.SizerFlags(1).Expand().DoubleBorder()
        )
        topSizer.Add(self.shapes["movie"][0], wx.SizerFlags(1).Expand().DoubleBorder())
        mainSizer.Add(topSizer, wx.SizerFlags(1).Expand())

        return mainSizer

    def SaveToPres(self, pres):
        slide = super().SaveToPres(pres)
        _ = self.shapes["textbox"][1].SaveToSlide(
            slide, Cm(0.5), Cm(4), Cm(16), Cm(2), scroll=True, font_size=12
        )

        movie, fs_slide, thn_img = self.shapes["movie"][0].SaveToSlide(
            pres, slide, Cm(33.87 / 2.0), Cm(0.37), Cm(8), Cm(6)
        )

        return slide, fs_slide


class wxStepSlide_wBOM(wxStepSlide):
    def __init__(
        self,
        parent,
        title,
        text,
        table_data1,
        table_data2,
        movie_file_name=None,
        movie_thumbnail_file_name=None,
        **kwargs,
    ):
        super().__init__(
            parent,
            title,
            text,
            movie_file_name=movie_file_name,
            movie_thumbnail_file_name=movie_thumbnail_file_name,
            **kwargs,
        )

        self.shapes["table"][0].SetFromDataframe(table_data1)
        self.shapes["table"][0].tableCtrl.AutoSizeColumns()
        self.shapes["table"][0].tableCtrl.AutoSizeRows()

        self.shapes["table"][1].SetFromDataframe(table_data2)
        self.shapes["table"][1].tableCtrl.AutoSizeColumns()
        self.shapes["table"][1].tableCtrl.AutoSizeRows()

    def MakeSlideLayout(self):
        mainSizer = super().MakeSlideLayout()

        self.AddTable(title="Components:")
        self.AddTable(title="Tool:")

        bottomSizer = wx.BoxSizer(wx.HORIZONTAL)
        bottomSizer.Add(
            self.shapes["table"][0], wx.SizerFlags(1).Expand().DoubleBorder()
        )
        bottomSizer.Add(
            self.shapes["table"][1], wx.SizerFlags(1).Expand().DoubleBorder()
        )
        mainSizer.Add(bottomSizer, wx.SizerFlags(1).Expand())

        return mainSizer

    def SaveToPres(self, pres):
        slide, fs_slide = super().SaveToPres(pres)

        _ = self.shapes["table"][0].SaveToSlide(
            slide, Cm(0.2), Cm(8), Cm(12), Cm(2.4), font_size=10
        )

        toolTable = self.shapes["table"][1].SaveToSlide(
            slide, Cm(12.4), Cm(8), Cm(12), Cm(2.4), font_size=10
        )
        tblf = toolTable._element.graphic.graphicData.tbl
        tblf[0][-1].text = "{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}"

        return slide, fs_slide


class wxPresentation(wx.Notebook):
    def __init__(self, parent, **kwargs):
        super().__init__(parent, **kwargs)
        self.AddPage(wxTitleSlide(self), "Setup")

    def __getitem__(self, idx):
        return self.GetPage(idx)

    def AddStepSlide(
        self, title, text, file_name, bom_tables=None, movie_thumbnail_file_name=None
    ):
        slideIdx = str(self.GetPageCount())
        if bom_tables is not None:
            self.AddPage(
                wxStepSlide_wBOM(
                    self,
                    title,
                    text,
                    bom_tables[0],
                    bom_tables[1],
                    file_name,
                    movie_thumbnail_file_name=movie_thumbnail_file_name,
                ),
                slideIdx,
            )
        else:
            self.AddPage(
                wxStepSlide(
                    self,
                    title,
                    text,
                    file_name,
                    movie_thumbnail_file_name=movie_thumbnail_file_name,
                ),
                slideIdx,
            )

    def Save(self, destination):
        pres = pptx.Presentation()
        hidden_slides = []
        for islide in range(self.GetPageCount()):
            slide_add = self.GetPage(islide).SaveToPres(pres)
            if isinstance(slide_add, tuple):
                hidden_slides.append(slide_add[1])

        for slide in hidden_slides:
            move_slide(pres, pres.slides.index(slide), -1)

        pres.save(destination)
