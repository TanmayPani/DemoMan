import wx
import sys
import pptx
import spacy
from os import path, listdir, makedirs, rename, environ, pathsep
import multiprocessing
import importlib.util
from multiprocessing import Process, set_start_method, freeze_support

# import imageio_ffmpeg
# ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()


def setup_ffmpeg():
    if getattr(sys, "frozen", False):
        base_path = sys._MEIPASS
    else:
        base_path = path.dirname(path.abspath(__file__))

    # The folder where you put ffmpeg.exe, ffprobe.exe, ffplay.exe
    ffmpeg_dir = path.join(base_path, "ffmpeg", "bin")
    current_path = environ.get("PATH", "")
    new_path = f"{current_path}{pathsep}{ffmpeg_dir}"
    environ["PATH"] = new_path


def get_model_path(mname):
    # If running in PyInstaller bundle
    if getattr(sys, "frozen", False):
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    # environ['KOKORO_MODEL_DIR'] = path.join(base_path, 'models')
    # Point to the folder containing your .pth files
    return os.path.join(base_path, mname)


setup_ffmpeg()

# def get_kmodel_path():
#     if hasattr(sys, '_MEIPASS'):
#         # Navigate to the folder containing config.cfg you found earlier
#         return os.path.join(sys._MEIPASS, "en_core_web_sm", "en_core_web_sm-3.8.0")
#     return "en_core_web_sm"
# model_pth = get_kmodel_path()

# def force_load_spacy_model():
#     if hasattr(sys, '_MEIPASS'):
#         # 1. Locate the actual model folder (ensure this matches your _internal structure)
#         model_dir = os.path.join(sys._MEIPASS, "en_core_web_sm", "en_core_web_sm-3.8.0")

#         # 2. Manually load it as a module named 'en_core_web_sm'
#         spec = importlib.util.spec_from_file_location("en_core_web_sm", os.path.join(model_dir, "__init__.py"))
#         model_module = importlib.util.module_from_spec(spec)
#         sys.modules["en_core_web_sm"] = model_module
#         spec.loader.exec_module(model_module)

#         return "en_core_web_sm" # Return the string name, which now works!
#     return "en_core_web_sm"


def load_spacy_model():
    if hasattr(sys, "_MEIPASS"):
        # Start at the root of your model collection
        base_search_path = os.path.join(sys._MEIPASS, "en_core_web_sm")

        # 1. Check if the config is right here
        if os.path.exists(os.path.join(base_search_path, "config.cfg")):
            return base_search_path

        # 2. Check one level deeper (for that 3.8.0 folder)
        if os.path.exists(base_search_path):
            for folder in os.listdir(base_search_path):
                subfolder = os.path.join(base_search_path, folder)
                if os.path.isdir(subfolder) and "config.cfg" in os.listdir(subfolder):
                    return subfolder

    return "en_core_web_sm"  # Dev fallback


# Use it before KPipeline
# model_pth = force_load_spacy_model()

# spaxypathtop = get_model_path("en_core_web_sm")
# spaxypath = os.path.join(spaxypathtop, "en_core_web_sm-3.8.0")
# print(spaxypath)

# # Force the pipeline to use your bundled file
# model_pth = get_res("kokoro-v1_0.pth")


# pipeline = KPipeline(lang_code='a', model=model_path)

# environ["PATH"] += pathsep + r'C:\Users\ramdi\KokoroEnv\ffmpeg-2026-03-01-git-862338fe31-full_build\ffmpeg-2026-03-01-git-862338fe31-full_build\bin'
# environ["IMAGEIO_FFMPEG_EXE"] = ffmpeg_exe
# environ["IMAGEIO_FFMPEG_EXE"] = r"C:\Users\ramdi\KokoroEnv\ffmpeg-2026-03-01-git-862338fe31-full_build\ffmpeg-2026-03-01-git-862338fe31-full_build\bin\ffmpeg.exe"


import torch

torch.serialization.add_safe_globals(
    ["omegaconf.listconfig.ListConfig", "omegaconf.dictconfig.DictConfig"]
)

import pickle
import shutil
from natsort import natsorted
import whisperx
import subprocess
import numpy as np
import faster_whisper
import ctranslate2
from re import findall, sub
import soundfile as sf
from pathlib import Path
from pptx.util import Pt
from pptx.util import Cm
from pptx.util import Inches
from transformers import Wav2Vec2ForCTC, Wav2Vec2Processor
from kokoro.pipeline import KPipeline
from threading import Thread
from pandas import read_excel, read_pickle, DataFrame, read_csv
from moviepy import concatenate_videoclips, concatenate_audioclips
from moviepy import VideoFileClip, AudioFileClip, CompositeAudioClip


class PageOne(wx.Panel):
    def __init__(self, parent):
        wx.Panel.__init__(self, parent)


class Thor(wx.Frame):
    def __init__(self, parent, id, nlp):
        wx.Frame.__init__(
            self,
            parent,
            id,
            "Thorlabs Instruction Transcription Interface",
            size=(1000, 600),
        )
        self.nlp = nlp

        # wx.Frame.__init__(self,parent,id,'Thorlabs Origin Interface',size=(750,500))

        panel = wx.Panel(self)
        self.nb = wx.Notebook(panel)

        sizer = wx.BoxSizer()
        sizer.Add(self.nb, 1, wx.EXPAND)
        panel.SetSizer(sizer)
        page1 = PageOne(self.nb)

        self.nb.AddPage(page1, "Assembly Presentation Generator")
        self.CorePath = None
        self.CorePathSelector = None
        self.StructurePath = None
        self.LCB = None
        self.VideoSliceButton = None
        self.VideoCombButton = None
        self.AudioPath = None
        self.VideoPath = None
        self.TextPath = None
        self.SegPath = None
        self.BOMWriterCB = None
        self.yoffset = 50
        Choices = ["Original", "Rewrite"]
        self.LCB = wx.StaticText(
            page1, -1, "Audio Transcription", pos=(40, 340 + self.yoffset)
        )
        self.LCB.Hide()
        self.AudioWriterCB = wx.ComboBox(
            page1,
            wx.ID_ANY,
            choices=Choices,
            style=wx.CB_READONLY,
            pos=(40, 360 + self.yoffset),
        )
        self.AudioWriterCB.SetValue("Original")
        self.Bind(wx.EVT_COMBOBOX, self.on_combo_selection, self.AudioWriterCB)
        self.AudioWriterCB.Hide()
        self.Layout()
        BOMChoices = ["No BOM", "BOM"]
        # self.LBW = wx.StaticText(page1,-1,'BOM Availability',pos=(40,10 +self.yoffset)); self.LCB.Hide()
        self.BOMWriterCB = wx.ComboBox(
            page1, wx.ID_ANY, choices=BOMChoices, style=wx.CB_READONLY, pos=(40, 10)
        )
        self.BOMWriterCB.SetValue("No BOM")
        self.Bind(wx.EVT_COMBOBOX, self.BOMSelection, self.BOMWriterCB)
        # self.Layout()

        self.WordSegments = None
        self.StepAudio = None
        self.StepVideo = None
        self.VideoClips = None
        self.ChangeLines = None
        self.VideoSlicerButton = None
        self.BKeys = None
        self.StepPIDs = None
        self.StepTools = None

        self.PPTXTitleSlide = None
        self.PPTXTitleSlideSubtitle = None
        self.s = wx.StaticText(self, label="Status: Idle", pos=(40, 250))

        status = self.CreateStatusBar()
        status_font = self.GetStatusBar().GetFont()

        new_font_size = status_font.GetPointSize() + 2  # Increase size by 2 points
        new_font = wx.Font(
            new_font_size,
            status_font.GetFamily(),
            status_font.GetStyle(),
            status_font.GetWeight(),
            status_font.GetUnderlined(),
            status_font.GetFaceName(),
        )
        self.GetStatusBar().SetFont(new_font)
        self.Show()

        # Displaying Text
        font = wx.Font(10, wx.DECORATIVE, wx.NORMAL, wx.BOLD)

        CPS = wx.Button(
            page1, label="Core Path", pos=(40, 10 + self.yoffset), size=(150, -1)
        )
        self.Bind(wx.EVT_BUTTON, self.PathSelector, CPS)
        VCB = wx.Button(
            page1, label="Video Compressor", pos=(40, 50 + self.yoffset), size=(150, -1)
        )
        self.Bind(wx.EVT_BUTTON, self.VideoCombination, VCB)
        VSB = wx.Button(
            page1, label="Video Slicer", pos=(40, 90 + self.yoffset), size=(150, -1)
        )
        #         self.Bind(wx.EVT_BUTTON,self.start_transcription,VSB)
        self.Bind(wx.EVT_BUTTON, self.TranscriptionModel, VSB)
        self.VideoSliceButton = VSB
        self.VideoCombButton = VCB
        self.L1 = wx.StaticText(page1, -1, "PPTX Name", pos=(270, 10 + self.yoffset))
        self.L2 = wx.StaticText(page1, -1, "Title Slide", pos=(270, 50 + self.yoffset))
        self.L3 = wx.StaticText(page1, -1, "Author", pos=(270, 90 + self.yoffset))
        self.PPTXName = wx.TextCtrl(page1, pos=(360, 10 + self.yoffset), size=(200, -1))
        self.PPTXName.SetValue("Test")
        self.PPTXTitleSlideTitle = wx.TextCtrl(
            page1, pos=(360, 50 + self.yoffset), size=(200, -1)
        )
        self.PPTXTitleSlideTitle.SetValue("Test")
        self.PPTXTitleSlideSubtitle = wx.TextCtrl(
            page1, pos=(360, 90 + self.yoffset), size=(200, -1)
        )
        self.PPTXTitleSlideSubtitle.SetValue("Test")

    def closebutton(self, event):
        self.Close(True)

    def closewindow(self, event):
        self.Destroy()

    def PathSelector(self, e):

        fileDialog = wx.DirDialog(
            frame, "Choose a directory", "", wx.DD_DEFAULT_STYLE | wx.DD_DIR_MUST_EXIST
        )
        if fileDialog.ShowModal() == wx.ID_OK:
            result = fileDialog.GetPath().replace("\\", "/") + "/"
            self.CorePath = result
            Structure = listdir(result)
            self.status = self.SetStatusText(
                "Successfully loaded: " + result + " -- Structure: " + str(Structure)
            )
        else:
            self.status = self.SetStatusText(
                "Command Terminated."
            )  # Untoggle the button
        #                 print(self.filepath)
        if "BOM" not in str(Structure):
            makedirs((result + "BOM/"), exist_ok=True)
            self.StructurePath = result + "BOM/"
            self.FileMover()
            self.StructurePath = None
        else:
            self.StructurePath = result + "BOM/"
            self.FileMover()
            self.StructurePath = None
        if "Videos" not in str(Structure):
            makedirs((result + "Videos/"), exist_ok=True)
            self.StructurePath = result + "Videos/"
            self.FileMover()
            self.StructurePath = None
        else:
            self.StructurePath = result + "Videos/"
            self.FileMover()
            self.StructurePath = None
        if "StepSegs" not in str(Structure):
            makedirs((result + "StepSegs/"), exist_ok=True)
        if "StepSegsAudio" not in str(Structure):
            makedirs((result + "StepSegsAudio/"), exist_ok=True)
        if "StepSegsTxt" not in str(Structure):
            makedirs((result + "StepSegsTxt/"), exist_ok=True)

        self.TextPath = self.CorePath + "StepSegsTxt/"
        self.VideoPath = self.CorePath + "Videos/"
        self.SegPath = self.CorePath + "StepSegs/"
        self.AudioPath = self.CorePath + "StepSegsAudio/"
        self.AudioWriterCB.Show()
        self.LCB.Show()
        self.Layout()

    def FileMover(self):
        source_path = self.CorePath
        destination_path = self.StructurePath
        if "BOM" in self.StructurePath:
            CoreFiles = [i for i in listdir(source_path) if ".xlsm" in i]
        if "Videos" in self.StructurePath:
            CoreFiles = [
                i for i in listdir(source_path) if ".MP4".casefold() in i.casefold()
            ]
        if len(CoreFiles) > 0:
            # try:
            # Move the files
            for i in CoreFiles:
                print(str(source_path) + "/" + str(i), str(destination_path))
                shutil.move(str(source_path) + "/" + str(i), str(destination_path))
                self.status = self.SetStatusText(
                    "Moved files into -- "
                    + str(destination_path)
                    + str(CoreFiles[0])
                    + " successfully."
                )  # Untoggle the button

        else:
            self.status = self.SetStatusText(
                "Nothing to move, or Error moving file into -- "
                + str(destination_path)
                + ". Verify it exists in core directory or if moved earlier."
            )  # Untoggle the button

    def VideoCombination(self, e):
        self.s.SetLabel("Begin video combination/renaming.")

        # logger = (self.statusbar)
        Corepath = self.CorePath
        # VideoPath=Corepath+'Videos/';self.VideoPath = VideoPath;self.SegPath = self.CorePath+'StepSegs/'#; self.AudioPath = Corepath+'StepSegsAudio/'
        VideoPath = self.VideoPath
        self.VideoCombButton.Disable()
        self.Layout()

        def work():
            try:
                if "combined.mp4" not in listdir(VideoPath):
                    VFiles = [
                        i
                        for i in listdir(VideoPath)
                        if ".MP4".casefold() in i.casefold()
                        if ".jpg" not in i
                    ]
                    # print(VFiles)
                    if len(VFiles) > 1:
                        wx.CallAfter(
                            self.s.SetLabel,
                            "Multiple videos identified. Initiate sorting sequence.",
                        )
                        SegVidInds = natsorted(VFiles)
                        #                 SegVidInds = np.argsort(np.concatenate([np.double(findall(r'\d+', i.replace('.MP4'.casefold(),'').replace(Corepath+"Videos/",''))) for i in VFiles]))
                        #                 SortedVFiles = [VideoFileClip(VideoPath+VFiles[SegVidInds[i]]) for i in range(len(SegVidInds))]
                        SortedVFiles = [
                            VideoFileClip(VideoPath + i) for i in SegVidInds
                        ]
                        SortedVFiles[0].save_frame(
                            self.CorePath + "StepSegs/" + "FirstFrame.jpg", t=0
                        )
                        wx.CallAfter(
                            self.s.SetLabel, "Sorting complete. Initiate concatenation."
                        )

                        CombinedVid = concatenate_videoclips(SortedVFiles)
                        wx.CallAfter(
                            self.s.SetLabel,
                            "Concatenation complete. Writing video to output. Please be patient.",
                        )

                        CombinedVid.write_videofile(
                            VideoPath + "combined.mp4",
                            temp_audiofile="temp-audio.m4a",
                            remove_temp=True,
                            audio_codec="aac",
                            codec="libx264",
                            threads=8,
                            logger=None,
                            preset="veryfast",
                        )
                        # CombinedVid.write_videofile(Corepath+'/Videos/combined.mp4',temp_audiofile="temp-audio.m4a", remove_temp=True, audio_codec="aac", codec="libx264",logger=logger)
                        # wx.CallAfter(self.statusbar.SetStatusText, "Export Complete!")
                        # video = VideoFileClip(VideoPath+'combined.mp4')
                        wx.CallAfter(
                            self.s.SetLabel,
                            "Video writing complete. Now writing audio to output.",
                        )

                        CombinedVid.audio.write_audiofile(
                            self.CorePath + "StepSegsAudio/" + "combined.mp3"
                        )
                        wx.CallAfter(
                            self.s.SetLabel,
                            "Audio successfully embedded in video. Please proceed to next step.",
                        )

                        #                         CombinedVid.audio.write_audiofile(self.CorePath+'StepSegsAudio/'+"wavf_subs.wav", fps=16000, nbytes=2, codec='pcm_s16le', logger=None)
                        # self.AudioPath = self.CorePath+'StepSegsAudio/'
                        CombinedVid.audio.close()
                        CombinedVid.close()
                    if len(VFiles) == 1:
                        wx.CallAfter(
                            self.s.SetLabel,
                            "Single video identified. Initiate renaming logic.",
                        )
                        rename(VideoPath + VFiles[0], VideoPath + "combined.mp4")
                        video = VideoFileClip(VideoPath + "combined.mp4")
                        video.save_frame(
                            self.CorePath + "StepSegs/" + "FirstFrame.jpg", t=0
                        )
                        wx.CallAfter(self.s.SetLabel, "Audio embedding start.")
                        video.audio.write_audiofile(
                            self.CorePath + "StepSegsAudio/" + "combined.mp3"
                        )
                        wx.CallAfter(
                            self.s.SetLabel,
                            "Audio embedding successful. Please proceed to next step",
                        )
                        #                 video.audio.write_audiofile(self.CorePath+'StepSegsAudio/'+"combined.wav", fps=16000, nbytes=2, codec='pcm_s16le', logger=None)
                        #                         CombinedVid.audio.write_audiofile(self.CorePath+'StepSegsAudio/'+"wavf_subs.wav", fps=16000, nbytes=2, codec='pcm_s16le', logger=None)
                        video.audio.close()
                        video.close()
                    if len(VFiles) == 0:
                        wx.CallAfter(self.s.SetLabel, "No videos avaialable in path :(")

                        # self.status = self.SetStatusText('No videos avaiable in path :(')  # Untoggle the button
                else:
                    video = VideoFileClip(VideoPath + "combined.mp4")
                    video.save_frame(
                        self.CorePath + "StepSegs/" + "FirstFrame.jpg", t=0
                    )
                    video.close()
                    wx.CallAfter(
                        self.s.SetLabel,
                        "combined.mp4 already exists. Please delete or move if another instance needs to be created.",
                    )

                    # self.status = self.SetStatusText('combined.mp4 already exists. Please delete or move if another instance needs to be created.')  # Untoggle the button
                    wx.CallAfter(
                        self.s.SetLabel,
                        "Rewriting audio of existing combined.mp4 file.",
                    )
                    video = VideoFileClip(VideoPath + "combined.mp4")
                    video.audio.write_audiofile(
                        self.CorePath + "StepSegsAudio/" + "combined.mp3"
                    )
                    wx.CallAfter(
                        self.s.SetLabel,
                        "Audio successfully embedded. Please proceed to next step.",
                    )

            except Exception as e:
                # wx.CallAfter(self.s.SetLabel, f"Error: {(AudioFile+'combined.mp3')}")
                wx.CallAfter(self.s.SetLabel, f"Error: {str(e)}")
                wx.CallAfter(self.VideoSliceButton.Enable)

        Thread(target=work, daemon=True).start()
        self.VideoCombButton.Enable()
        self.Layout()

    def TimeSlicer(self):
        word_segments = self.WordSegments
        StartSeg = -1
        EndSeg = -1
        Step = []
        for i in range(len(word_segments)):
            if i + 2 < len(word_segments):
                if (
                    "start" in word_segments[i]["word"].casefold()
                    and "step" in word_segments[i + 1]["word"].casefold()
                ):
                    StartSeg = i
                if (
                    "end" in word_segments[i]["word"].casefold()
                    or "and" in word_segments[i]["word"].casefold()
                    or "finish" in word_segments[i]["word"].casefold()
                    or "stop" in word_segments[i]["word"].casefold()
                ) and "step" in word_segments[i + 1]["word"].casefold():
                    EndSeg = i + 3
            if StartSeg != -1 and EndSeg != -1:
                Step.append(word_segments[StartSeg:EndSeg])
                StartSeg = -1
                EndSeg = -1
        StichedStep = [" ".join([j["word"] for j in i]) for i in Step]
        # print('End Stitching. Attempting any corrections')
        for i in StichedStep:
            if "and step" in i.casefold() or "And step" in i.casefold():
                i.casefold().replace("and step", "end step")
        # print('Corrections successful!')
        CompleteStepTiming = [
            [Step[i][0]["start"], Step[i][len(Step[i]) - 1]["end"]]
            for i in range(len(Step))
        ]
        return Step, StichedStep, CompleteStepTiming

    def TranscriptionModel(self, e):
        self.status = self.SetStatusText(
            "Transcription sequence initiated."
        )  # Untoggle the button
        #         AudioFile = self.AudioPath+'combined.mp3'
        # self.AudioPath = self.CorePath+'StepSegsAudio/';
        AudioFile = self.AudioPath
        #         if torch.cuda.is_available():
        #             device = "cuda"
        #         else:
        #             device='cpu'
        device = "cpu"
        batch_size = 8
        if AudioFile == "":
            raise Exception("Please enter valid filepath")
        # if batch_size <16:
        #     compute_type = "int8"
        # else:
        #     compute_type = "float16"
        compute_type = "int8"
        self.VideoSliceButton.Disable()
        self.Layout()

        self.status = self.SetStatusText("Initiating thread...")  # Untoggle the button

        self.s.SetLabel("Loading & Transcribing... " + str(device))

        def work():
            try:
                AudioFile = self.AudioPath
                device = "cpu"
                compute_type = "int8"
                batch_size = 1
                wx.CallAfter(self.s.SetLabel, "Attempting to load large-v3 model.")
                Model = whisperx.load_model(
                    "large-v3", device, compute_type=compute_type
                )
                wx.CallAfter(self.s.SetLabel, "Attempting to load audio.")
                LoadedAudio = whisperx.load_audio(Path(AudioFile + "combined.mp3"))
                wx.CallAfter(self.s.SetLabel, "Initiate transcription.")
                PreliminaryTranscription = Model.transcribe(
                    LoadedAudio, batch_size=batch_size, language="en"
                )
                wx.CallAfter(self.s.SetLabel, "Time projection.")
                AlignModel, Metadata = whisperx.load_align_model(
                    language_code="en", device=device
                )
                AlignedTranscription = whisperx.align(
                    PreliminaryTranscription["segments"],
                    AlignModel,
                    Metadata,
                    AudioFile + "combined.mp3",
                    device,
                    return_char_alignments=False,
                )
                wx.CallAfter(self.s.SetLabel, "Begin writing.")
                self.WordSegments = AlignedTranscription["word_segments"]
                wx.CallAfter(
                    self.s.SetLabel,
                    "Alignment sequence complete. Initiate step isolation sequence.",
                )
                self.FullSteps, self.StichedSteps, self.StichedTiming = (
                    self.TimeSlicer()
                )
                wx.CallAfter(
                    self.s.SetLabel,
                    "Step Isolation sequence complete. Obtaining audio slices",
                )
                self.StepAudio = [
                    self.AudioWriter(
                        self.StichedSteps[i],
                        self.AudioPath + "AFHeart",
                        i,
                        self.AudioWriterCB.GetValue(),
                    )
                    for i in range(len(self.StichedSteps))
                ]
                wx.CallAfter(
                    self.s.SetLabel, "Audio slices obtained. Initiate video rendering."
                )
                self.StepVideo = [
                    self.VideoStepWriter(self.StichedTiming[i], i)
                    for i in range(len(self.StichedTiming))
                ]
                wx.CallAfter(
                    self.s.SetLabel,
                    "Videos successfully rendered. Presentation creation initiated.",
                )
                if self.BOMWriterCB.GetValue() == "BOM":
                    wx.CallAfter(self.s.SetLabel, "Extracting BOM data.")
                    self.BOMWriter()
                    wx.CallAfter(
                        self.s.SetLabel, "BOM data extracted. Starting presentation."
                    )
                    self.PresentationGenerator()
                else:
                    wx.CallAfter(self.s.SetLabel, "Starting presentation.")
                    self.NoBOMPresentationGenerator()
                wx.CallAfter(
                    self.s.SetLabel,
                    "KokoroAI Processes Complete. Thank you for your patience!",
                )  # Thread-safe update

            except Exception as e:
                # wx.CallAfter(self.s.SetLabel, f"Error: {(AudioFile+'combined.mp3')}")
                wx.CallAfter(self.s.SetLabel, f"Error: {str(e)}")
                wx.CallAfter(self.VideoSliceButton.Enable)

        # 2. Fire it off in a thread immediately
        Thread(target=work, daemon=True).start()
        # self.s.SetLabel("Processing...")
        #         self.StepVideo = [self.VideoStepWriter(self.StichedTiming[i],i) for i in range(len(self.StichedTiming))]
        #         self.status = self.SetStatusText('Excel Gener.')  # Untoggle the button
        self.VideoSliceButton.Enable()
        self.Layout()

    def AudioWriter(self, StepSegments, Tag, Index, ManualOverride):
        pipeline = self.nlp
        if ManualOverride == "Original":
            generator = pipeline(
                StepSegments,
                voice="af_heart",  # <= change voice here
                speed=0.8,
                split_pattern=r"\n+",
            )
        else:
            data = read_csv(self.TextPath + "Corrected_Files.csv", encoding="latin1")
            Changelines = data["0"].to_list()
            tmp = [
                sub(r" ?\[.*?\]", "", item).replace("-", " dash ")
                for item in Changelines
            ]
            FText = [sub(r"(\d)", r"\1 ", i) for i in tmp]
            generator = pipeline(
                FText[Index],
                voice="af_heart",  # <= change voice here
                #             self.ChangeLines[Index], voice='af_heart', # <= change voice here
                #                 ManualOverride, voice='af_heart', # <= change voice here
                speed=0.8,
                split_pattern=r"\n+",
            )
        AudioClips = []
        TextClips = []
        for i, (gs, ps, audio) in enumerate(generator):
            AudioClips.append(Tag + "Step" + str(i) + ".mp3")
            sf.write(
                Tag + "Step" + str(i) + ".mp3", audio, 24000
            )  # save each audio file
        FullAudio = concatenate_audioclips([AudioFileClip(i) for i in AudioClips])
        FullAudio.write_audiofile(Tag + "FullStep" + str(Index) + ".mp3", 24000)
        return Tag + "FullStep" + str(Index) + ".mp3"

    def VideoStepWriter(self, e, Index):
        # print('Start Video Step Writing')
        VideoPath = self.VideoPath
        AudioPath = self.AudioPath
        Timing = self.StichedTiming
        VideoClips = []
        #         NoSubVideo = VideoFileClip(VideoPath+'combined.mp4').without_audio()
        Video = VideoFileClip(VideoPath + "combined.mp4").without_audio()
        #         Video = CompositeVideoClip([NoSubVideo, self.subs])

        MiniAud = AudioFileClip(self.StepAudio[Index])
        MiniVid = Video[Timing[Index][0] : Timing[Index][1]]
        MiniVid = MiniVid.with_audio(MiniAud)
        #         print('Video Sliced')
        #         print(self.SegPath+"AFHeart"+str(Index)+".mp4")
        MiniVid.write_videofile(
            self.SegPath + "AFHeart" + str(Index) + ".mp4",
            codec="libx264",
            audio_codec="aac",
            preset="veryfast",
            logger=None,
            threads=8,
        )

        VideoClips.append(self.SegPath + "AFHeart" + str(Index) + ".mp4")
        self.VideoClips = VideoClips

    def on_complete(self, segments):
        self.s.SetLabel(f"Done! {len(segments)} segments found.")
        self.VideoSliceButton.Enable()

    #     def on_start_export(self, event):
    #         self.VideoSliceButton.Disable() # Prevent multiple clicks
    #         thread = threading.Thread(target=self.render_comb_video)
    #         thread.start()
    def render_comb_video(self, clip, output_path, threads):
        clip.write_videofile(
            output_path,
            temp_audiofile="temp-audio.m4a",
            remove_temp=True,
            audio_codec="aac",
            codec="libx264",
            threads=threads,
            preset="ultrafast",
            logger=None,
        )

    #         wx.CallAfter(self.on_export_finished)
    #     def on_export_finished(self):
    #         self.VideoSliceButton.Enable()
    #         wx.MessageBox("Export complete!", "Info")
    def BOMWriter(self):
        self.BOMPath = self.CorePath + "BOM/"
        BomPath = self.BOMPath
        self.BKeys, self.StepPIDs, self.StepTools = self.standardized_excel_reader(
            BomPath
        )
        self.ComponentWriter(
            self.TextPath, "AFHeartTxt.pkl", [self.StepPIDs, self.StepTools]
        )

    def ComponentWriter(self, fpath, name, comps):
        with open(fpath + name, "wb") as f:
            pickle.dump(comps, f)

    def ComponentReader(self, fpath, name):
        with open(fpath + name, "rb") as f:
            comps = pickle.load(f)
        return comps

    def standardized_excel_reader(self, fpath):
        Files = [i for i in listdir(fpath) if ".xlsm" in i or ".xlsx" in i]
        # print(Files)
        FullSheet = read_excel(fpath + Files[0], sheet_name=None)
        SheetData = FullSheet.items()
        keys, dfs = [[] for i in range(2)]
        for key, df in SheetData:
            keys.append(key)
            dfs.append(df)

        StepKeys, StepPIDs, StepTools = [[] for i in range(3)]
        count = 0
        for i in dfs:
            if len(i["Item number"]) > 0:
                mask = i["Item number"] == "Tool"
                i["grouper"] = mask.cumsum()
                ig = {
                    group_key: group_df for group_key, group_df in i.groupby("grouper")
                }
                ig[1].columns = ig[1].iloc[0]
                ig[1] = ig[1].loc[:, :"Quantity"]
                StepKeys.append(keys[count])
                StepPIDs.append(
                    ig[0].dropna().reset_index(drop=True).drop(["grouper"], axis=1)
                )
                StepTools.append(
                    ig[1]["Tool Description"].dropna().reset_index(drop=True)[1:]
                )

            count += 1
        return keys, StepPIDs, StepTools

    def PresentationVideoWriter(self, vid, pres, stepindex, addslide):
        left = Cm(33.87 / 2)
        top = Cm(0.37)
        width = Cm(8)
        height = Cm(6)
        mime_type = "video/mp4"

        if addslide == True:
            sf = pres.slide_layouts[5]  # Get a title only slide layout
            s = pres.slides.add_slide(sf)
            Title1 = s.shapes.title
            s.text = "Step " + str(stepindex)
            FFPath = self.CorePath + "StepSegs/" + "FirstFrame.jpg"
            movie_shape = s.shapes.add_movie(
                vid,
                left,
                top,
                width,
                height,
                mime_type=mime_type,
                poster_frame_image=FFPath,
            )
        else:
            FFPath = self.CorePath + "StepSegs/" + "FirstFrame.jpg"
            movie_shape = pres.slides[stepindex].shapes.add_movie(
                vid,
                left,
                top,
                width,
                height,
                mime_type=mime_type,
                poster_frame_image=FFPath,
            )

    def PresentationGenerator(self):
        Subtitle = self.PPTXTitleSlideTitle.GetValue()
        CorePath = self.PPTXTitleSlideSubtitle.GetValue()
        PPTX_name = self.PPTXName.GetValue()
        # Create a presentation object (or load an existing one)
        # StepData = read_pickle(CorePath+'StepSegsTxt/'+'StepData.pkl')
        #         if self.AudioWriterCB.GetValue() == 'Original':
        StepData = read_pickle(self.TextPath + "AFHeartTxt.pkl")
        #         print('Original')
        #         else:
        #             print(read_pickle(self.TextPath+'AFHeartTxt.pkl'))
        #             print(self.ChangeLines)
        #         StepData = self.ChangeLines
        pres = pptx.Presentation()
        # blank_slide_layout = pres.slide_layouts[6] # Get a blank slide layout
        title_slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(title_slide_layout)
        Title1 = slide.shapes.title
        Title1.text = "Assembly Instruction Model"
        subtitle = slide.placeholders[1]
        subtitle.text = Subtitle
        sf2 = pres.slide_layouts[5]  # Get a title only slide layout

        SegInds = np.argsort(
            np.concatenate(
                [
                    np.double(
                        findall(
                            r"\d+", i.replace(".mp4", "").replace(self.SegPath, "")[-3:]
                        )
                    )
                    for i in listdir(self.SegPath)
                    if "AFHeart" in i
                ]
            )
        )
        StepVids = listdir(self.SegPath)
        FullStepVidPaths = [self.SegPath + i for i in StepVids if ".jpg" not in i]
        StichedDF = DataFrame(self.StichedSteps).to_csv(
            self.TextPath + "Uncorrected_Steps.csv"
        )
        VidParam = False
        offset = 1
        for i in range(len(StepData[0])):
            if len(StepData[0]) == len(FullStepVidPaths):
                self.TableWriter(StepData[0][i], pres, "Component", i + offset)
                self.TableWriter(StepData[1][i], pres, "Tool", i + offset)
                self.PresentationVideoWriter(
                    FullStepVidPaths[SegInds[i]], pres, i + offset, False
                )
                if self.AudioWriterCB.GetValue() == "Original":
                    self.TextWriter(self.StichedSteps[i], pres, i + offset)
                else:
                    data = read_csv(
                        self.TextPath + "Corrected_Files.csv", encoding="latin1"
                    )
                    Changelines = data["0"].to_list()
                    #                     self.ChangeLines = [sub(r" ?\[.*?\]","", item) for item in Changelines]
                    self.ChangeLines = Changelines
                    # self.status = self.SetStatusText('Rewrite requested. Please make sure to modify a file named Corrected_Files.csv, with the lines only occupying the same number of cells as the Uncorrected_Steps.csv file does.')
                    self.TextWriter(self.ChangeLines[i], pres, i + offset)
            else:
                self.TableWriter(StepData[0][i], pres, "Component", i + offset)
                self.TableWriter(StepData[1][i], pres, "Tool", i + offset)
                #                 self.TextWriter(self.StichedSteps[i],pres,i+offset)
                if self.AudioWriterCB.GetValue() == "Original":
                    self.TextWriter(self.StichedSteps[i], pres, i + offset)
                else:
                    data = read_csv(
                        self.TextPath + "Corrected_Files.csv", encoding="latin1"
                    )
                    Changelines = data["0"].to_list()
                    #                     self.ChangeLines = [sub(r" ?\[.*?\]","", item) for item in Changelines]
                    self.ChangeLines = Changelines
                    self.TextWriter(self.ChangeLines[i], pres, i + offset)
                VidParam = True
        if VidParam == True:
            for i in range(len(FullStepVidPaths)):
                self.PresentationVideoWriter(
                    FullStepVidPaths[SegInds[i]], pres, i + offset, False
                )

        pres.save(self.CorePath + PPTX_name + ".pptx")

    def NoBOMPresentationGenerator(self):
        MainTitle = self.PPTXTitleSlideTitle.GetValue()
        Subtitle = self.PPTXTitleSlideSubtitle.GetValue()
        PPTX_name = self.PPTXName.GetValue()
        pres = pptx.Presentation()
        title_slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(title_slide_layout)
        Title1 = slide.shapes.title
        Title1.text = MainTitle
        subtitle = slide.placeholders[1]
        subtitle.text = Subtitle

        SegInds = np.argsort(
            np.concatenate(
                [
                    np.double(
                        findall(
                            r"\d+", i.replace(".mp4", "").replace(self.SegPath, "")[-3:]
                        )
                    )
                    for i in listdir(self.SegPath)
                    if "AFHeart" in i
                ]
            )
        )

        StepVids = listdir(self.SegPath)
        FullStepVidPaths = [self.SegPath + i for i in StepVids if ".jpg" not in i]
        offset = 1
        StichedDF = DataFrame(self.StichedSteps).to_csv(
            self.TextPath + "Uncorrected_Steps.csv"
        )
        if self.AudioWriterCB.GetValue() == "Original":
            #             print('OG, if')
            #             print(FullStepVidPaths)
            #             print(SegInds)
            for i in range(len(FullStepVidPaths)):
                print(SegInds[i])
                self.PresentationVideoWriter(
                    FullStepVidPaths[SegInds[i]], pres, i + offset, True
                )
                self.TextWriter(self.StichedSteps[i], pres, i + offset)
        else:
            data = read_csv(self.TextPath + "Corrected_Files.csv", encoding="latin1")
            Changelines = data["0"].to_list()
            #             self.ChangeLines = [sub(r" ?\[.*?\]","", item) for item in Changelines]
            self.ChangeLines = Changelines
            for i in range(len(FullStepVidPaths)):
                self.PresentationVideoWriter(
                    FullStepVidPaths[SegInds[i]], pres, i + offset, True
                )
                self.TextWriter(self.ChangeLines[i], pres, i + offset)
        pres.save(self.CorePath + PPTX_name + ".pptx")

    def set_table_font_size(self, table, size_pt):
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(size_pt)

    def TableWriter(self, data, pres, tabletype, stepindex):
        if tabletype == "Component":
            rows = data.shape[0] + 1
            cols = data.shape[1]
            title_only_slide_layout = pres.slide_layouts[5]
            slide = pres.slides.add_slide(title_only_slide_layout)
            shapes = slide.shapes
            shapes.title.text = "Step " + str(stepindex)

            # rows = cols = 2
            left = Cm(0.2)
            top = Cm(8)
            width = Cm(12)
            height = Cm(0.8) * cols

            table = shapes.add_table(rows, cols, left, top, width, height).table
            for j, column_name in enumerate(data.columns):
                table.cell(0, j).text = column_name

            # Add data rows
            for i, row_data in enumerate(data.values.tolist()):
                for j, cell_data in enumerate(row_data):
                    table.cell(i + 1, j).text = str(cell_data)
            self.set_table_font_size(table, 10)
        if tabletype == "Tool":
            rows = data.shape[0] + 1
            cols = 1  # data.shape[1]
            data = DataFrame(list(data), columns=["Tools"])
            slide = pres.slides[stepindex]
            shapes = slide.shapes
            shapes.title.text = "Step " + str(stepindex)
            # shapes.title.text = 'Step '+str(stepindex-1)
            left = Cm(12.4)
            top = Cm(8)
            width = Cm(12)
            height = Cm(0.8) * cols
            tbl = shapes.add_table(rows, cols, left, top, width, height)
            table = tbl.table
            # table.style ='Medium Grid 1 Accent 2' #(9DCAF9ED-07DC-4A11-8D7F-57B35C25682E)
            # Access the underlying XML element of the table shape
            tblf = tbl._element.graphic.graphicData.tbl
            style_id = "{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}"
            # Inject the style ID into the XML
            tblf[0][-1].text = style_id

            for j, column_name in enumerate(data.columns):
                table.cell(0, j).text = column_name

            # Add data rows
            for i, row_data in enumerate(data.values.tolist()):
                for j, cell_data in enumerate(row_data):
                    table.cell(i + 1, j).text = str(cell_data)
        self.set_table_font_size(table, 10)

    def TextWriter(self, txt, pres, stepindex):
        slide = pres.slides[stepindex]
        left = Cm(0.5)
        top = Cm(4)
        width = Cm(16)
        height = Cm(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = txt
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

    def on_combo_selection(self, event):
        # 4. Get the selected value
        selected_item = event.GetEventObject().GetValue()
        if selected_item == "Original":
            self.status = self.SetStatusText("Raw audio transcription being used")
        else:
            data = read_csv(self.TextPath + "Corrected_Files.csv", encoding="latin1")
            Changelines = data["0"].to_list()
            #             print(Changelines)
            self.ChangeLines = data["0"].to_list()
            #             self.ChangeLines = [sub(r" ?\[.*?\]","", item) for item in Changelines]

            #             print(self.ChangeLines)
            self.status = self.SetStatusText(
                "Rewrite requested. Please make sure to modify a file named Corrected_Files.csv, with the lines only occupying the same number of cells as the Uncorrected_Steps.csv file does."
            )

    def BOMSelection(self, e):
        # 4. Get the selected value
        selected_item = e.GetEventObject().GetValue()
        return selected_item


if __name__ == "__main__":
    # REQUIRED: Must set 'spawn' before creating any processes (default on Windows)
    multiprocessing.freeze_support()

    try:
        set_start_method("spawn", force=True)
    except RuntimeError:
        pass
    try:
        # model_pth = force_load_spacy_model()
        model_pth = load_spacy_model()
        pipeline = KPipeline(lang_code="a", model=model_pth)
        print("Model and Pipeline loaded successfully!")
        # nlp_model = load_spacy_model()
        # print("Model loaded successfully!")
    except Exception as e:
        print(f"Failed to load model: {e}")
        nlp_model = None
    app = wx.App()
    frame = Thor(parent=None, id=-1, nlp=pipeline)
    frame.Show()
    app.MainLoop()
